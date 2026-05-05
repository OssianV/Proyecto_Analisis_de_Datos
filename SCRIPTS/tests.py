from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


# Colores
COLOR_FONDO = "F7F7F2"
COLOR_TEXTO = "1F2937"
COLOR_GRID = "D0D7DE"

COLOR_PRINCIPAL = "1D4E89"
COLOR_SECUNDARIO = "2A9D8F"
COLOR_NEUTRO = "8D99AE"


def rgb(hex_color):
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def agregar_barra_superior(slide, titulo, color):
    barra = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
    barra.fill.solid()
    barra.fill.fore_color.rgb = rgb(color)
    barra.line.color.rgb = rgb(color)

    caja_titulo = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(10.5), Inches(0.3))
    p = caja_titulo.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = titulo
    run.font.name = "Aptos"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = rgb(COLOR_FONDO)


def agregar_recuadro_blanco(slide):
    recuadro = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1.0), Inches(1.4),
        Inches(11.3), Inches(4.9)
    )
    recuadro.fill.solid()
    recuadro.fill.fore_color.rgb = RGBColor(255, 255, 255)
    recuadro.line.color.rgb = rgb(COLOR_GRID)
    return recuadro


def agregar_texto(slide, left, top, width, height, texto, tamano, color, negrita=False):
    caja = slide.shapes.add_textbox(left, top, width, height)
    tf = caja.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.name = "Aptos"
    run.font.size = Pt(tamano)
    run.font.bold = negrita
    run.font.color.rgb = rgb(color)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: portada
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLOR_FONDO)

    agregar_barra_superior(slide, "Presentacion", COLOR_PRINCIPAL)

    agregar_texto(
        slide,
        Inches(1.2), Inches(2.0),
        Inches(10.8), Inches(0.8),
        "Analisis de datos para decisiones futbolisticas",
        24,
        COLOR_TEXTO,
        True
    )

    agregar_texto(
        slide,
        Inches(1.2), Inches(3.0),
        Inches(10.8), Inches(0.8),
        "Scouting, perfiles de jugador y apoyo tactico",
        16,
        COLOR_NEUTRO,
        False
    )

    # Slide 2: grafica de ejemplo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLOR_FONDO)

    agregar_barra_superior(slide, "Analisis 01: Grafica", COLOR_PRINCIPAL)
    agregar_recuadro_blanco(slide)

    agregar_texto(
        slide,
        Inches(3.7), Inches(3.45),
        Inches(5.0), Inches(0.4),
        "Aqui ira la grafica",
        24,
        COLOR_NEUTRO,
        True
    )

    # Slide 3: analisis de ejemplo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLOR_FONDO)

    agregar_barra_superior(slide, "Lectura del analisis 01", COLOR_SECUNDARIO)
    agregar_recuadro_blanco(slide)

    agregar_texto(
        slide,
        Inches(1.4), Inches(1.9),
        Inches(10.5), Inches(3.7),
        "En este recuadro ira el texto del analisis.\n\nPuedes pegar aqui 2 o 3 parrafos completos, como los que escribiste en tu libreta.\n\nLa idea es que esta slide tenga exactamente el mismo tamano de recuadro que la slide de la grafica.",
        18,
        COLOR_TEXTO,
        False
    )

    # Slide 4: integrantes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLOR_FONDO)

    agregar_barra_superior(slide, "Integrantes", COLOR_PRINCIPAL)
    agregar_recuadro_blanco(slide)

    agregar_texto(
        slide,
        Inches(1.4), Inches(2.0),
        Inches(5.0), Inches(2.5),
        "Nombre Apellido 1\n\nNombre Apellido 2\n\nNombre Apellido 3\n\nNombre Apellido 4",
        20,
        COLOR_TEXTO,
        False
    )

    output_dir = Path(__file__).resolve().parents[1] / "OUTPUTS"
    output_dir.mkdir(exist_ok=True)

    output_file = output_dir / "sample_presentacion_simple.pptx"
    prs.save(output_file)

    print(f"Sample generado en: {output_file}")


if __name__ == "__main__":
    main()